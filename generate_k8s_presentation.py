from pptx import Presentation
from pptx.util import Inches, Pt

# Creamos la presentación
prs = Presentation()

# Ajustes para títulos y contenidos
title_slide_layout = prs.slide_layouts[0]  # Título
bullet_slide_layout = prs.slide_layouts[1] # Título + Cuerpo con bullets

##############################
# 1. Title Slide
##############################
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Kubernetes Concepts: An Advanced Overview for DevOps"
subtitle = slide.placeholders[1]
subtitle.text = "Presenter: [Your Name / Team]\nDate: [Date]\nOrganization: [Company/Team Name]"

##############################
# Función útil para crear slides con bullets
##############################
def create_bullet_slide(prs, title, bullet_points):
    """
    Crea una slide con un título y una lista de bullet points.
    prs: objeto Presentation
    title: str - texto para el título
    bullet_points: list - lista de strings (cada string será un bullet)
    """
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = title
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    # Evita que se duplique el primer bullet vacío
    tf.clear()
    
    for bp in bullet_points:
        p = tf.add_paragraph()
        p.text = bp
        p.level = 0  # Nivel de bullet

##############################
# 2. Agenda
##############################
agenda_points = [
    "1. Introduction & High-Level Overview",
    "2. Kubernetes Architecture",
    "3. The Kubernetes API",
    "4. Workloads (Pods, Controllers, Jobs, CronJobs)",
    "5. Services, Load Balancing & Networking",
    "6. Configuration (ConfigMaps, Secrets)",
    "7. Storage (Volumes, PV, PVC, StorageClass)",
    "8. Security & Policies",
    "9. Cluster Administration & Extension",
    "10. Q&A"
]
create_bullet_slide(prs, "Agenda", agenda_points)

##############################
# 3. Introduction & High-Level Overview
##############################
intro_points = [
    "Kubernetes is an open-source system for automating deployment, scaling, and management of containerized applications.",
    "Originally designed by Google, now maintained by CNCF.",
    "Solves challenges around container orchestration: scheduling, scaling, networking, and lifecycle management."
]
create_bullet_slide(prs, "What is Kubernetes?", intro_points)

key_concepts_points = [
    "Clusters: group of nodes (physical or virtual machines).",
    "Control Plane (Master Components): API server, etcd, Scheduler, Controller Manager.",
    "Nodes: run containerized applications (Pods).",
    "Declarative Model: desired state via YAML/JSON manifests.",
    "Idempotency: repeated application of config results in the same final state."
]
create_bullet_slide(prs, "Key Concepts Overview", key_concepts_points)

##############################
# 4. Kubernetes Architecture
##############################
control_plane_points = [
    "kube-apiserver: Front-end for the Kubernetes control plane.",
    "etcd: Key-value store for cluster data.",
    "kube-scheduler: Schedules pods to nodes based on requirements.",
    "kube-controller-manager: Runs various controller processes."
]
create_bullet_slide(prs, "Control Plane Components", control_plane_points)

node_components_points = [
    "kubelet: Agent on each node, ensures containers run correctly.",
    "kube-proxy: Maintains network rules for node, implements Service abstraction.",
    "Container Runtime: e.g. Docker, containerd, CRI-O."
]
create_bullet_slide(prs, "Node Components", node_components_points)

addons_points = [
    "CoreDNS: Internal DNS resolution for Services.",
    "Ingress Controllers: Manage external access.",
    "Metrics Server: Resource usage aggregator for autoscaling.",
    "Dashboard: Web UI for basic cluster management (use with caution in production)."
]
create_bullet_slide(prs, "Add-On Components", addons_points)

##############################
# 5. The Kubernetes API
##############################
api_model_points = [
    "RESTful API: All cluster operations exposed via kube-apiserver.",
    "API Groups & Versions: e.g., apps/v1, batch/v1.",
    "Declarative vs. Imperative: YAML manifests used in a declarative approach.",
    "Watch Mechanism: real-time streaming of changes to objects."
]
create_bullet_slide(prs, "API Object Model", api_model_points)

working_api_points = [
    "kubectl: CLI tool for Kubernetes (CRUD operations).",
    "Controllers: continuously reconcile actual state with desired state.",
    "Admission Controllers: validate or mutate requests before persisting."
]
create_bullet_slide(prs, "Working with the API", working_api_points)

##############################
# 6. Workloads
##############################
pods_points = [
    "Definition: A Pod is one or more containers with shared storage/network.",
    "Ephemeral nature: Pods are managed by controllers for scaling.",
    "Lifecycle phases: Pending → Running → Succeeded/Failed → Unknown.",
    "Multi-Container Pods: Common pattern is sidecar container."
]
create_bullet_slide(prs, "Pods: The Smallest Deployable Unit", pods_points)

controllers_points = [
    "Purpose: Ensure desired # of Pods, handle rolling updates, maintain state.",
    "Types: Deployments, ReplicaSets, StatefulSets, DaemonSets, Jobs, CronJobs."
]
create_bullet_slide(prs, "Overview of Controllers", controllers_points)

replicaset_points = [
    "Ensures a specific number of identical Pods are running.",
    "Usually used by Deployments rather than directly.",
    "Selector-based for matching Pods."
]
create_bullet_slide(prs, "ReplicaSet", replicaset_points)

deployment_points = [
    "Primary controller for stateless services.",
    "Rolling Updates & Rollbacks minimize downtime.",
    "Encapsulates ReplicaSet management for easier updates.",
    "Update Strategies: RollingUpdate (default) or Recreate."
]
create_bullet_slide(prs, "Deployment", deployment_points)

statefulset_points = [
    "Manages stateful apps needing stable network identity/storage.",
    "Stable Pod identifiers across rescheduling.",
    "Ordered deployment and scaling."
]
create_bullet_slide(prs, "StatefulSet", statefulset_points)

daemonset_points = [
    "Ensures one Pod per node (or subset of nodes).",
    "Useful for cluster-wide tasks: logging, monitoring, etc.",
    "Pods added automatically as new nodes join."
]
create_bullet_slide(prs, "DaemonSet", daemonset_points)

job_points = [
    "For batch or one-off tasks to completion.",
    "Ensures specified number of successful completions.",
    "Reschedules Pods on failure if needed."
]
create_bullet_slide(prs, "Job", job_points)

cronjob_points = [
    "Creates Jobs on a time-based schedule (cron format).",
    "Used for periodic tasks like backups, cleanup, etc."
]
create_bullet_slide(prs, "CronJob", cronjob_points)

##############################
# 7. Services, Load Balancing & Networking
##############################
services_points = [
    "Stable network endpoint for a set of Pods.",
    "Label selectors define which Pods receive traffic.",
    "Types: ClusterIP, NodePort, LoadBalancer, ExternalName."
]
create_bullet_slide(prs, "Services", services_points)

ingress_points = [
    "Manages external (HTTP/HTTPS) access to Services.",
    "Requires an Ingress Controller (Nginx, HAProxy, Traefik...).",
    "Supports routing rules, TLS termination, etc."
]
create_bullet_slide(prs, "Ingress", ingress_points)

dns_points = [
    "CoreDNS by default for internal name resolution.",
    "Resolves Service names to cluster IPs.",
    "Service FQDN format: <service>.<namespace>.svc.cluster.local"
]
create_bullet_slide(prs, "Cluster DNS", dns_points)

netpol_points = [
    "NetworkPolicy defines firewall rules at Pod level.",
    "Restricts ingress/egress based on labels/selectors.",
    "Requires CNI plugin supporting NetworkPolicy (e.g. Calico)."
]
create_bullet_slide(prs, "Network Policies", netpol_points)

##############################
# 8. Configuration
##############################
configmap_points = [
    "Stores non-confidential config data in key-value pairs.",
    "Decouples configuration from container images.",
    "Mounted as env vars or files."
]
create_bullet_slide(prs, "ConfigMaps", configmap_points)

secret_points = [
    "Stores sensitive data (passwords, tokens) base64-encoded.",
    "Similar usage as ConfigMaps, but for confidential values.",
    "Best practice: integrate with external secret management (Vault)."
]
create_bullet_slide(prs, "Secrets", secret_points)

##############################
# 9. Storage
##############################
volumes_points = [
    "Ephemeral volumes (like emptyDir) tied to Pod lifecycle.",
    "Containers in a Pod can share Volumes.",
    "Different from Docker volumes model; more ephemeral nature."
]
create_bullet_slide(prs, "Volumes", volumes_points)

pv_points = [
    "Abstract persistent storage from underlying providers.",
    "Lifecycle independent from the Pod.",
    "Reclaim Policy: Retain, Recycle, or Delete."
]
create_bullet_slide(prs, "Persistent Volumes (PV)", pv_points)

pvc_points = [
    "Requests storage (size, access modes, storage class).",
    "Auto-binds to matching PV or triggers dynamic provisioning.",
    "Access Modes: ReadWriteOnce, ReadOnlyMany, ReadWriteMany."
]
create_bullet_slide(prs, "Persistent Volume Claims (PVC)", pvc_points)

storageclass_points = [
    "Defines dynamic provisioning of storage.",
    "Specifies underlying storage type, parameters, etc.",
    "Different classes for SSD, HDD, region/zone, etc."
]
create_bullet_slide(prs, "StorageClass", storageclass_points)

##############################
# 10. Security
##############################
authn_points = [
    "Methods: x509 certs, bearer tokens, external (OIDC).",
    "Service Accounts for pods to securely interact with API."
]
create_bullet_slide(prs, "Authentication", authn_points)

authz_points = [
    "Modes: RBAC, ABAC, Webhook.",
    "Roles & ClusterRoles define permissions.",
    "RoleBinding & ClusterRoleBinding grant permissions."
]
create_bullet_slide(prs, "Authorization", authz_points)

pod_security_points = [
    "Pod Security Standards (PSS): Privileged, Baseline, Restricted.",
    "PodSecurityPolicy (PSP) deprecated in 1.25, replaced by Pod Security admission.",
    "Restricts privileged access, volume types, Linux capabilities."
]
create_bullet_slide(prs, "Pod Security", pod_security_points)

##############################
# 11. Policies
##############################
resource_quota_points = [
    "Limits resource consumption (CPU, memory, storage) at namespace level.",
    "Prevents a single team/app from monopolizing resources."
]
create_bullet_slide(prs, "Resource Quotas", resource_quota_points)

limit_range_points = [
    "Defines default requests/limits for containers in a namespace.",
    "Ensures min and max resource constraints per container."
]
create_bullet_slide(prs, "Limit Ranges", limit_range_points)

pdb_points = [
    "PodDisruptionBudget limits voluntary disruptions (node maintenance).",
    "Ensures minimum number or percentage of pods remain available."
]
create_bullet_slide(prs, "PodDisruptionBudget (PDB)", pdb_points)

##############################
# 12. Cluster Administration
##############################
cluster_admin_points = [
    "Lifecycle management: upgrades, backups, node operations.",
    "Logging & Monitoring: Metrics Server, Prometheus, Grafana.",
    "Autoscaling: HPA, VPA.",
    "High Availability: multiple control-plane instances, distributed etcd."
]
create_bullet_slide(prs, "Cluster Administration Essentials", cluster_admin_points)

##############################
# 13. Extending Kubernetes
##############################
extending_points = [
    "Custom Resource Definitions (CRD) to extend the API.",
    "Operators: advanced controllers managing CRDs.",
    "Examples: Database operators, message queue operators."
]
create_bullet_slide(prs, "Custom Resources & Operators", extending_points)

webhooks_points = [
    "Mutating & Validating Webhooks: custom logic on object creation/update.",
    "API Aggregation: combine multiple APIs behind Kubernetes APIServer."
]
create_bullet_slide(prs, "Admission Webhooks & API Aggregation", webhooks_points)

##############################
# 14. Summary & Next Steps
##############################
takeaways_points = [
    "Modular & extensible Kubernetes architecture.",
    "Pods as the base unit; controllers provide the operational power.",
    "Services & Ingress handle cluster-internal and external communication.",
    "Security & Policies crucial for multi-tenant environments.",
    "Extensibility allows custom workflows via CRDs & Operators."
]
create_bullet_slide(prs, "Key Takeaways", takeaways_points)

next_steps_points = [
    "Explore advanced topics: Operators, Service Mesh, GitOps, multi-cluster.",
    "Review official documentation on kubernetes.io/docs.",
    "Hands-on practice with Minikube or kind.",
    "Integrate with CI/CD for automated tests & deployments."
]
create_bullet_slide(prs, "Next Steps", next_steps_points)

##############################
# 15. Q&A
##############################
qanda_points = [
    "Open for discussion.",
    "Share real-world experiences and challenges."
]
create_bullet_slide(prs, "Q&A", qanda_points)

# Guardamos la presentación con nombre "Kubernetes_Advanced_DevOps.pptx"
prs.save("Kubernetes_Advanced_DevOps.pptx")

print("Fichero 'Kubernetes_Advanced_DevOps.pptx' generado correctamente.")

